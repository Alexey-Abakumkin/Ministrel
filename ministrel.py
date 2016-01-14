import argparse
from math import floor, exp

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR


PRESENTATION_WIDTH = 13004800 # 1024px
PRESENTATION_HEIGHT = 9753600 # 768px


def create_presentation(output, texts, master_slide, font_color):

    # Here is exponential regression to fit character without hyphenation
    def estimate_font_size(str_length):
        a = 31.52805
        b = 266.4258
        c = 0.091539
        pt = floor(a + b*exp(-c*str_length))
        if pt < 32:
            return 32
        if pt > 200:
            return 200
        return pt

    def estimate_max_length(text):
        return max([len(x) for x in text])

    def prepare_text(text):
        return '\n'.join(text)

    def decorate_shape(shape, text, font_color):
        run = shape.text_frame.paragraphs[0].add_run()
        run.font.size = Pt(estimate_font_size(estimate_max_length(text)))
        run.text = prepare_text(text)

        if font_color == 'white':
            color = RGBColor(255, 255, 255)
        else:
            color = RGBColor(0, 0, 0)

        run.font.color.rgb = color

    def create_slide(slide, text, font_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PRESENTATION_WIDTH, PRESENTATION_HEIGHT)
        shape.fill.background()
        shape.line.color.rgb = RGBColor(0, 0, 0)
        decorate_shape(shape, text, font_color)

    p = Presentation(master_slide)
    for i, text in enumerate(texts):
        if i == 0:
            slide = p.slides[0]
        else:
            slide = p.slides.add_slide(p.slide_layouts[11])
        create_slide(slide, text, font_color)

    p.save(output)


def read_song(filename):
    with open(filename) as f:
        song = f.readlines()

    texts = "".join(song).split("\n\n")
    return [x.split('\n') for x in texts]


if __name__ == '__main__':
    def get_app_args():
        parser = argparse.ArgumentParser(
            description='The program creates pptx presentation from text file with lyrics'
        )

        parser.add_argument('-m', '--master',
                            dest='slide',
                            default='master.pptx',
                            help='destination to master pptx slide (default: "master.pptx")')

        parser.add_argument('-c', '--font-color',
                            dest='color',
                            default='white',
                            choices=['white', 'black'],
                            help='color of the presentation font (default: "white")')

        parser.add_argument('input_file',
                            metavar='INPUT',
                            help='text file name contains lyrics of the song')

        parser.add_argument('output_file',
                            metavar='OUTPUT',
                            help='output presentation destination')

        return parser.parse_args()

    args = get_app_args()

    texts = read_song(args.input_file)
    create_presentation(args.output_file, texts, args.slide, args.color)

    print('Presentation "{}" was created'.format(args.output_file))
